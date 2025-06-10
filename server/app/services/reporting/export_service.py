"""
Export Service

Provides comprehensive export capabilities:
- PDF reports with custom layouts and styling
- Excel workbooks with multiple sheets and formatting
- CSV files with various delimiters and encodings
- PowerPoint presentations with charts and data
- Word documents with tables and formatting
- JSON and XML data exports
"""

import json
import csv
import io
import os
import tempfile
import uuid
from datetime import datetime
from typing import Dict, List, Any, Optional, BinaryIO
import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DocxInches
import xml.etree.ElementTree as ET


class ExportService:
    """Service for exporting reports in various formats"""

    def __init__(self):
        self.temp_dir = tempfile.mkdtemp()
        self.supported_formats = {
            'pdf': {'name': 'PDF Document', 'mime_type': 'application/pdf', 'extension': '.pdf'},
            'excel': {'name': 'Excel Workbook', 'mime_type': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'extension': '.xlsx'},
            'csv': {'name': 'CSV File', 'mime_type': 'text/csv', 'extension': '.csv'},
            'powerpoint': {'name': 'PowerPoint Presentation', 'mime_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation', 'extension': '.pptx'},
            'word': {'name': 'Word Document', 'mime_type': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'extension': '.docx'},
            'json': {'name': 'JSON File', 'mime_type': 'application/json', 'extension': '.json'},
            'xml': {'name': 'XML File', 'mime_type': 'application/xml', 'extension': '.xml'}
        }

    def export_report(self, report_data: Dict[str, Any], format_type: str, 
                     export_config: Dict[str, Any] = None) -> Dict[str, Any]:
        """Export report in the specified format"""
        
        if format_type not in self.supported_formats:
            raise ValueError(f"Unsupported export format: {format_type}")
        
        export_config = export_config or {}
        
        try:
            if format_type == 'pdf':
                result = self._export_pdf(report_data, export_config)
            elif format_type == 'excel':
                result = self._export_excel(report_data, export_config)
            elif format_type == 'csv':
                result = self._export_csv(report_data, export_config)
            elif format_type == 'powerpoint':
                result = self._export_powerpoint(report_data, export_config)
            elif format_type == 'word':
                result = self._export_word(report_data, export_config)
            elif format_type == 'json':
                result = self._export_json(report_data, export_config)
            elif format_type == 'xml':
                result = self._export_xml(report_data, export_config)
            else:
                raise ValueError(f"Export format {format_type} not implemented")
            
            return {
                'success': True,
                'format': format_type,
                'file_path': result['file_path'],
                'file_name': result['file_name'],
                'file_size': result['file_size'],
                'mime_type': self.supported_formats[format_type]['mime_type'],
                'export_time': datetime.utcnow().isoformat(),
                'record_count': len(report_data.get('data', []))
            }
            
        except Exception as e:
            return {
                'success': False,
                'format': format_type,
                'error': str(e)
            }

    def _export_pdf(self, report_data: Dict[str, Any], config: Dict[str, Any]) -> Dict[str, Any]:
        """Export report as PDF"""
        
        file_name = config.get('file_name', f'report_{uuid.uuid4().hex[:8]}.pdf')
        file_path = os.path.join(self.temp_dir, file_name)
        
        # Create PDF document
        doc = SimpleDocTemplate(
            file_path,
            pagesize=config.get('page_size', A4),
            leftMargin=config.get('left_margin', 0.75) * inch,
            rightMargin=config.get('right_margin', 0.75) * inch,
            topMargin=config.get('top_margin', 1) * inch,
            bottomMargin=config.get('bottom_margin', 1) * inch
        )
        
        # Build content
        story = []
        styles = getSampleStyleSheet()
        
        # Title
        title = config.get('title', 'Report')
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=config.get('title_font_size', 18),
            textColor=colors.toColor(config.get('title_color', '#333333')),
            spaceAfter=20,
            alignment=1  # Center alignment
        )
        story.append(Paragraph(title, title_style))
        
        # Metadata
        metadata = report_data.get('metadata', {})
        if config.get('include_metadata', True) and metadata:
            story.append(Spacer(1, 12))
            
            meta_text = f"Generated: {metadata.get('generated_at', datetime.utcnow().isoformat())}<br/>"
            meta_text += f"Records: {len(report_data.get('data', []))}<br/>"
            if metadata.get('filter_count', 0) > 0:
                meta_text += f"Filters Applied: {metadata['filter_count']}<br/>"
            
            story.append(Paragraph(meta_text, styles['Normal']))
            story.append(Spacer(1, 20))
        
        # Data table
        data = report_data.get('data', [])
        if data and config.get('include_data_table', True):
            # Prepare table data
            if config.get('fields'):
                headers = [field.get('alias', field.get('field', '')) for field in config['fields']]
                table_data = [headers]
                for row in data:
                    table_row = [str(row.get(field.get('alias', field.get('field', '')), '')) for field in config['fields']]
                    table_data.append(table_row)
            else:
                # Use all fields from first row
                headers = list(data[0].keys()) if data else []
                table_data = [headers]
                for row in data:
                    table_row = [str(row.get(col, '')) for col in headers]
                    table_data.append(table_row)
            
            # Limit rows if specified
            max_rows = config.get('max_rows', 1000)
            if len(table_data) > max_rows + 1:  # +1 for header
                table_data = table_data[:max_rows + 1]
                story.append(Paragraph(f"<i>Note: Only first {max_rows} rows shown</i>", styles['Normal']))
                story.append(Spacer(1, 12))
            
            # Create table
            table = Table(table_data)
            
            # Table styling
            table_style = [
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), config.get('header_font_size', 10)),
                ('FONTSIZE', (0, 1), (-1, -1), config.get('data_font_size', 8)),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]
            
            # Alternate row colors
            if config.get('alternate_row_colors', True):
                for i in range(1, len(table_data)):
                    if i % 2 == 0:
                        table_style.append(('BACKGROUND', (0, i), (-1, i), colors.lightgrey))
            
            table.setStyle(TableStyle(table_style))
            story.append(table)
        
        # Charts (if provided)
        charts = config.get('charts', [])
        for chart_config in charts:
            if 'image_data' in chart_config:
                story.append(Spacer(1, 20))
                # Save chart image to temporary file
                chart_file = os.path.join(self.temp_dir, f'chart_{uuid.uuid4().hex[:8]}.png')
                with open(chart_file, 'wb') as f:
                    f.write(chart_config['image_data'])
                
                # Add chart to PDF
                img = Image(chart_file)
                img.drawHeight = config.get('chart_height', 3) * inch
                img.drawWidth = config.get('chart_width', 6) * inch
                story.append(img)
        
        # Build PDF
        doc.build(story)
        
        file_size = os.path.getsize(file_path)
        
        return {
            'file_path': file_path,
            'file_name': file_name,
            'file_size': file_size
        }

    def _export_excel(self, report_data: Dict[str, Any], config: Dict[str, Any]) -> Dict[str, Any]:
        """Export report as Excel workbook"""
        
        file_name = config.get('file_name', f'report_{uuid.uuid4().hex[:8]}.xlsx')
        file_path = os.path.join(self.temp_dir, file_name)
        
        # Create workbook
        wb = Workbook()
        
        # Remove default sheet
        wb.remove(wb.active)
        
        # Data sheet
        data_sheet = wb.create_sheet("Data")
        data = report_data.get('data', [])
        
        if data:
            # Headers
            if config.get('fields'):
                headers = [field.get('alias', field.get('field', '')) for field in config['fields']]
            else:
                headers = list(data[0].keys()) if data else []
            
            # Write headers
            for col_idx, header in enumerate(headers, 1):
                cell = data_sheet.cell(row=1, column=col_idx, value=header)
                cell.font = Font(bold=True)
                cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
                cell.alignment = Alignment(horizontal='center')
                cell.border = Border(
                    left=Side(style='thin'),
                    right=Side(style='thin'),
                    top=Side(style='thin'),
                    bottom=Side(style='thin')
                )
            
            # Write data
            for row_idx, row in enumerate(data, 2):
                for col_idx, header in enumerate(headers, 1):
                    value = row.get(header, '')
                    cell = data_sheet.cell(row=row_idx, column=col_idx, value=value)
                    cell.border = Border(
                        left=Side(style='thin'),
                        right=Side(style='thin'),
                        top=Side(style='thin'),
                        bottom=Side(style='thin')
                    )
                    
                    # Apply zebra striping
                    if config.get('zebra_striping', True) and row_idx % 2 == 0:
                        cell.fill = PatternFill(start_color="F2F2F2", end_color="F2F2F2", fill_type="solid")
            
            # Auto-adjust column widths
            if config.get('auto_fit_columns', True):
                for column in data_sheet.columns:
                    max_length = 0
                    column_letter = column[0].column_letter
                    for cell in column:
                        try:
                            if len(str(cell.value)) > max_length:
                                max_length = len(str(cell.value))
                        except:
                            pass
                    adjusted_width = min(max_length + 2, 50)
                    data_sheet.column_dimensions[column_letter].width = adjusted_width
        
        # Summary sheet
        if config.get('include_summary', True):
            summary_sheet = wb.create_sheet("Summary")
            
            # Report metadata
            summary_data = [
                ["Report Title", config.get('title', 'Report')],
                ["Generated", datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')],
                ["Total Records", len(data)],
                ["Columns", len(headers) if data else 0]
            ]
            
            metadata = report_data.get('metadata', {})
            if metadata.get('filter_count', 0) > 0:
                summary_data.append(["Filters Applied", metadata['filter_count']])
            if metadata.get('data_sources'):
                summary_data.append(["Data Sources", ', '.join(metadata['data_sources'])])
            
            # Write summary data
            for row_idx, (label, value) in enumerate(summary_data, 1):
                summary_sheet.cell(row=row_idx, column=1, value=label).font = Font(bold=True)
                summary_sheet.cell(row=row_idx, column=2, value=value)
            
            # Auto-adjust column widths
            summary_sheet.column_dimensions['A'].width = 20
            summary_sheet.column_dimensions['B'].width = 30
        
        # Charts sheet
        charts_config = config.get('charts', [])
        if charts_config and data:
            charts_sheet = wb.create_sheet("Charts")
            
            # Create sample charts based on data
            for i, chart_config in enumerate(charts_config):
                chart_type = chart_config.get('type', 'bar')
                
                if chart_type == 'bar' and len(headers) >= 2:
                    chart = BarChart()
                    chart.title = chart_config.get('title', f'Chart {i+1}')
                    
                    # Add data (first 10 rows max for performance)
                    data_range = Reference(data_sheet, min_col=2, min_row=1, max_row=min(11, len(data)+1), max_col=2)
                    categories = Reference(data_sheet, min_col=1, min_row=2, max_row=min(11, len(data)+1))
                    
                    chart.add_data(data_range, titles_from_data=True)
                    chart.set_categories(categories)
                    
                    charts_sheet.add_chart(chart, f"A{i*15 + 1}")
                    
                elif chart_type == 'line' and len(headers) >= 2:
                    chart = LineChart()
                    chart.title = chart_config.get('title', f'Chart {i+1}')
                    
                    data_range = Reference(data_sheet, min_col=2, min_row=1, max_row=min(11, len(data)+1), max_col=2)
                    categories = Reference(data_sheet, min_col=1, min_row=2, max_row=min(11, len(data)+1))
                    
                    chart.add_data(data_range, titles_from_data=True)
                    chart.set_categories(categories)
                    
                    charts_sheet.add_chart(chart, f"A{i*15 + 1}")
                    
                elif chart_type == 'pie' and len(headers) >= 2:
                    chart = PieChart()
                    chart.title = chart_config.get('title', f'Chart {i+1}')
                    
                    data_range = Reference(data_sheet, min_col=2, min_row=1, max_row=min(11, len(data)+1), max_col=2)
                    categories = Reference(data_sheet, min_col=1, min_row=2, max_row=min(11, len(data)+1))
                    
                    chart.add_data(data_range, titles_from_data=True)
                    chart.set_categories(categories)
                    
                    charts_sheet.add_chart(chart, f"A{i*15 + 1}")
        
        # Save workbook
        wb.save(file_path)
        file_size = os.path.getsize(file_path)
        
        return {
            'file_path': file_path,
            'file_name': file_name,
            'file_size': file_size
        }

    def _export_csv(self, report_data: Dict[str, Any], config: Dict[str, Any]) -> Dict[str, Any]:
        """Export report as CSV"""
        
        file_name = config.get('file_name', f'report_{uuid.uuid4().hex[:8]}.csv')
        file_path = os.path.join(self.temp_dir, file_name)
        
        data = report_data.get('data', [])
        
        if not data:
            # Create empty CSV with headers if specified
            with open(file_path, 'w', newline='', encoding=config.get('encoding', 'utf-8')) as csvfile:
                writer = csv.writer(
                    csvfile, 
                    delimiter=config.get('delimiter', ','),
                    quotechar=config.get('quote_char', '"'),
                    quoting=csv.QUOTE_MINIMAL
                )
                if config.get('fields'):
                    headers = [field.get('alias', field.get('field', '')) for field in config['fields']]
                    writer.writerow(headers)
        else:
            # Determine headers
            if config.get('fields'):
                headers = [field.get('alias', field.get('field', '')) for field in config['fields']]
            else:
                headers = list(data[0].keys())
            
            with open(file_path, 'w', newline='', encoding=config.get('encoding', 'utf-8')) as csvfile:
                writer = csv.writer(
                    csvfile, 
                    delimiter=config.get('delimiter', ','),
                    quotechar=config.get('quote_char', '"'),
                    quoting=csv.QUOTE_MINIMAL
                )
                
                # Write headers
                if config.get('include_headers', True):
                    writer.writerow(headers)
                
                # Write data
                for row in data:
                    csv_row = []
                    for header in headers:
                        value = row.get(header, '')
                        # Handle None values
                        if value is None:
                            value = config.get('null_value', '')
                        csv_row.append(str(value))
                    writer.writerow(csv_row)
        
        file_size = os.path.getsize(file_path)
        
        return {
            'file_path': file_path,
            'file_name': file_name,
            'file_size': file_size
        }

    def _export_powerpoint(self, report_data: Dict[str, Any], config: Dict[str, Any]) -> Dict[str, Any]:
        """Export report as PowerPoint presentation"""
        
        file_name = config.get('file_name', f'report_{uuid.uuid4().hex[:8]}.pptx')
        file_path = os.path.join(self.temp_dir, file_name)
        
        # Create presentation
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = config.get('title', 'Report')
        subtitle.text = f"Generated on {datetime.utcnow().strftime('%Y-%m-%d')}\n{len(report_data.get('data', []))} records"
        
        # Summary slide
        if config.get('include_summary', True):
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = 'Report Summary'
            
            tf = body_shape.text_frame
            tf.text = f'Total Records: {len(report_data.get("data", []))}'
            
            metadata = report_data.get('metadata', {})
            if metadata.get('data_sources'):
                p = tf.add_paragraph()
                p.text = f'Data Sources: {", ".join(metadata["data_sources"])}'
            
            if metadata.get('filter_count', 0) > 0:
                p = tf.add_paragraph()
                p.text = f'Filters Applied: {metadata["filter_count"]}'
        
        # Data slide
        data = report_data.get('data', [])
        if data and config.get('include_data_table', True):
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add title
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = 'Data Table'
            title_frame.paragraphs[0].font.size = Inches(0.3)
            title_frame.paragraphs[0].font.bold = True
            
            # Prepare table data
            if config.get('fields'):
                headers = [field.get('alias', field.get('field', '')) for field in config['fields']]
            else:
                headers = list(data[0].keys()) if data else []
            
            # Limit data for PowerPoint
            max_rows = config.get('max_rows', 10)
            limited_data = data[:max_rows]
            
            # Create table
            rows = min(len(limited_data) + 1, 15)  # +1 for header, max 15 rows
            cols = min(len(headers), 8)  # Max 8 columns
            
            table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table
            
            # Set column widths
            for i in range(cols):
                table.columns[i].width = Inches(9 / cols)
            
            # Add headers
            for col_idx, header in enumerate(headers[:cols]):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                cell.text_frame.paragraphs[0].font.bold = True
            
            # Add data
            for row_idx, row in enumerate(limited_data[:rows-1], 1):
                for col_idx, header in enumerate(headers[:cols]):
                    cell = table.cell(row_idx, col_idx)
                    value = row.get(header, '')
                    cell.text = str(value)[:50]  # Limit cell content length
        
        # Charts slides
        charts_config = config.get('charts', [])
        for chart_config in charts_config:
            if 'image_data' in chart_config:
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Add chart title
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                title_frame = title_shape.text_frame
                title_frame.text = chart_config.get('title', 'Chart')
                title_frame.paragraphs[0].font.size = Inches(0.3)
                title_frame.paragraphs[0].font.bold = True
                
                # Save chart image to temporary file
                chart_file = os.path.join(self.temp_dir, f'chart_{uuid.uuid4().hex[:8]}.png')
                with open(chart_file, 'wb') as f:
                    f.write(chart_config['image_data'])
                
                # Add chart image
                slide.shapes.add_picture(chart_file, Inches(1), Inches(2), Inches(8), Inches(5))
        
        # Save presentation
        prs.save(file_path)
        file_size = os.path.getsize(file_path)
        
        return {
            'file_path': file_path,
            'file_name': file_name,
            'file_size': file_size
        }

    def _export_word(self, report_data: Dict[str, Any], config: Dict[str, Any]) -> Dict[str, Any]:
        """Export report as Word document"""
        
        file_name = config.get('file_name', f'report_{uuid.uuid4().hex[:8]}.docx')
        file_path = os.path.join(self.temp_dir, file_name)
        
        # Create document
        doc = Document()
        
        # Title
        title = doc.add_heading(config.get('title', 'Report'), 0)
        title.alignment = 1  # Center alignment
        
        # Metadata
        metadata = report_data.get('metadata', {})
        if config.get('include_metadata', True):
            doc.add_paragraph(f"Generated: {metadata.get('generated_at', datetime.utcnow().isoformat())}")
            doc.add_paragraph(f"Records: {len(report_data.get('data', []))}")
            
            if metadata.get('filter_count', 0) > 0:
                doc.add_paragraph(f"Filters Applied: {metadata['filter_count']}")
            
            if metadata.get('data_sources'):
                doc.add_paragraph(f"Data Sources: {', '.join(metadata['data_sources'])}")
            
            doc.add_paragraph("")  # Empty line
        
        # Data table
        data = report_data.get('data', [])
        if data and config.get('include_data_table', True):
            doc.add_heading('Data', level=1)
            
            # Prepare table data
            if config.get('fields'):
                headers = [field.get('alias', field.get('field', '')) for field in config['fields']]
            else:
                headers = list(data[0].keys()) if data else []
            
            # Limit rows for Word document
            max_rows = config.get('max_rows', 100)
            limited_data = data[:max_rows]
            
            # Create table
            table = doc.add_table(rows=1, cols=len(headers))
            table.style = 'Table Grid'
            
            # Add headers
            hdr_cells = table.rows[0].cells
            for i, header in enumerate(headers):
                hdr_cells[i].text = str(header)
                # Make header bold
                for paragraph in hdr_cells[i].paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
            
            # Add data rows
            for row in limited_data:
                row_cells = table.add_row().cells
                for i, header in enumerate(headers):
                    value = row.get(header, '')
                    row_cells[i].text = str(value)
            
            if len(data) > max_rows:
                doc.add_paragraph(f"Note: Only first {max_rows} rows shown")
        
        # Charts
        charts_config = config.get('charts', [])
        if charts_config:
            doc.add_heading('Charts', level=1)
            
            for chart_config in charts_config:
                if 'image_data' in chart_config:
                    # Save chart image to temporary file
                    chart_file = os.path.join(self.temp_dir, f'chart_{uuid.uuid4().hex[:8]}.png')
                    with open(chart_file, 'wb') as f:
                        f.write(chart_config['image_data'])
                    
                    # Add chart title
                    if chart_config.get('title'):
                        doc.add_heading(chart_config['title'], level=2)
                    
                    # Add chart image
                    doc.add_picture(chart_file, width=DocxInches(6))
                    doc.add_paragraph("")  # Empty line
        
        # Save document
        doc.save(file_path)
        file_size = os.path.getsize(file_path)
        
        return {
            'file_path': file_path,
            'file_name': file_name,
            'file_size': file_size
        }

    def _export_json(self, report_data: Dict[str, Any], config: Dict[str, Any]) -> Dict[str, Any]:
        """Export report as JSON"""
        
        file_name = config.get('file_name', f'report_{uuid.uuid4().hex[:8]}.json')
        file_path = os.path.join(self.temp_dir, file_name)
        
        # Prepare export data
        export_data = {
            'metadata': {
                'title': config.get('title', 'Report'),
                'generated_at': datetime.utcnow().isoformat(),
                'record_count': len(report_data.get('data', [])),
                'export_format': 'json'
            }
        }
        
        # Add original metadata
        if report_data.get('metadata'):
            export_data['metadata'].update(report_data['metadata'])
        
        # Add data
        if config.get('include_data', True):
            export_data['data'] = report_data.get('data', [])
        
        # Add fields information
        if config.get('fields'):
            export_data['fields'] = config['fields']
        
        # Write JSON file
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(
                export_data, 
                f, 
                indent=config.get('indent', 2),
                ensure_ascii=config.get('ensure_ascii', False),
                default=str  # Handle datetime and other non-serializable objects
            )
        
        file_size = os.path.getsize(file_path)
        
        return {
            'file_path': file_path,
            'file_name': file_name,
            'file_size': file_size
        }

    def _export_xml(self, report_data: Dict[str, Any], config: Dict[str, Any]) -> Dict[str, Any]:
        """Export report as XML"""
        
        file_name = config.get('file_name', f'report_{uuid.uuid4().hex[:8]}.xml')
        file_path = os.path.join(self.temp_dir, file_name)
        
        # Create root element
        root = ET.Element('report')
        
        # Metadata
        metadata_elem = ET.SubElement(root, 'metadata')
        ET.SubElement(metadata_elem, 'title').text = config.get('title', 'Report')
        ET.SubElement(metadata_elem, 'generated_at').text = datetime.utcnow().isoformat()
        ET.SubElement(metadata_elem, 'record_count').text = str(len(report_data.get('data', [])))
        ET.SubElement(metadata_elem, 'export_format').text = 'xml'
        
        # Add original metadata
        original_metadata = report_data.get('metadata', {})
        for key, value in original_metadata.items():
            ET.SubElement(metadata_elem, key).text = str(value)
        
        # Data
        if config.get('include_data', True):
            data_elem = ET.SubElement(root, 'data')
            
            for row in report_data.get('data', []):
                record_elem = ET.SubElement(data_elem, 'record')
                for key, value in row.items():
                    # Clean key name for XML
                    clean_key = ''.join(c for c in key if c.isalnum() or c == '_')
                    if clean_key and not clean_key[0].isdigit():
                        elem = ET.SubElement(record_elem, clean_key)
                        elem.text = str(value) if value is not None else ''
                    else:
                        # Use generic field name for invalid keys
                        elem = ET.SubElement(record_elem, 'field')
                        elem.set('name', key)
                        elem.text = str(value) if value is not None else ''
        
        # Fields information
        if config.get('fields'):
            fields_elem = ET.SubElement(root, 'fields')
            for field in config['fields']:
                field_elem = ET.SubElement(fields_elem, 'field')
                for key, value in field.items():
                    field_elem.set(key, str(value))
        
        # Write XML file
        tree = ET.ElementTree(root)
        tree.write(
            file_path, 
            encoding='utf-8', 
            xml_declaration=True,
            method='xml'
        )
        
        # Pretty print if requested
        if config.get('pretty_print', True):
            self._pretty_print_xml(file_path)
        
        file_size = os.path.getsize(file_path)
        
        return {
            'file_path': file_path,
            'file_name': file_name,
            'file_size': file_size
        }

    def _pretty_print_xml(self, file_path: str):
        """Pretty print XML file"""
        try:
            import xml.dom.minidom
            
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            dom = xml.dom.minidom.parseString(content)
            pretty_content = dom.toprettyxml(indent='  ', encoding='utf-8')
            
            with open(file_path, 'wb') as f:
                f.write(pretty_content)
                
        except Exception:
            # If pretty printing fails, keep original
            pass

    def get_supported_formats(self) -> Dict[str, Dict]:
        """Return supported export formats"""
        return self.supported_formats

    def validate_export_config(self, format_type: str, 
                             config: Dict[str, Any]) -> Dict[str, Any]:
        """Validate export configuration"""
        
        errors = []
        warnings = []
        
        if format_type not in self.supported_formats:
            errors.append(f"Unsupported export format: {format_type}")
            return {'is_valid': False, 'errors': errors, 'warnings': warnings}
        
        # Format-specific validations
        if format_type == 'csv':
            delimiter = config.get('delimiter', ',')
            if len(delimiter) != 1:
                errors.append("CSV delimiter must be a single character")
            
            encoding = config.get('encoding', 'utf-8')
            valid_encodings = ['utf-8', 'utf-16', 'latin-1', 'ascii']
            if encoding not in valid_encodings:
                warnings.append(f"Encoding '{encoding}' may not be supported")
        
        elif format_type == 'pdf':
            max_rows = config.get('max_rows', 1000)
            if max_rows > 5000:
                warnings.append(f"Large number of rows ({max_rows}) may impact PDF generation performance")
        
        elif format_type == 'excel':
            if config.get('charts') and len(config['charts']) > 10:
                warnings.append("Large number of charts may impact Excel file size")
        
        # File name validation
        file_name = config.get('file_name', '')
        if file_name:
            invalid_chars = ['<', '>', ':', '"', '|', '?', '*', '\\', '/']
            if any(char in file_name for char in invalid_chars):
                errors.append("File name contains invalid characters")
        
        return {
            'is_valid': len(errors) == 0,
            'errors': errors,
            'warnings': warnings
        }

    def cleanup_temp_files(self, older_than_hours: int = 24):
        """Clean up temporary export files"""
        
        import time
        
        now = time.time()
        cutoff = now - (older_than_hours * 3600)
        
        cleaned_count = 0
        
        try:
            for file_name in os.listdir(self.temp_dir):
                file_path = os.path.join(self.temp_dir, file_name)
                if os.path.isfile(file_path):
                    file_mtime = os.path.getmtime(file_path)
                    if file_mtime < cutoff:
                        os.remove(file_path)
                        cleaned_count += 1
        except Exception:
            pass  # Ignore cleanup errors
        
        return cleaned_count